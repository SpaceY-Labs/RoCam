#!/usr/bin/env python3
"""Convert the HTML presentation to PPTX by screenshotting each slide."""

import os
import pathlib
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

SCRIPT_DIR = pathlib.Path(__file__).parent.resolve()
HTML_PATH = SCRIPT_DIR / "presentation.html"
OUTPUT_PATH = SCRIPT_DIR / "RoCam_Final_Presentation.pptx"
SLIDES_DIR = SCRIPT_DIR / "slide_images"
SLIDE_W, SLIDE_H = 1920, 1080
SCALE = 5  # 5x = 9600×5400 maximum quality

def capture_slides():
    SLIDES_DIR.mkdir(exist_ok=True)
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": SLIDE_W, "height": SLIDE_H},
            device_scale_factor=SCALE,
        )
        page.goto(f"file://{HTML_PATH}", wait_until="networkidle")
        page.wait_for_timeout(3000)

        slides = page.query_selector_all(".slide")
        print(f"Found {len(slides)} slides (scale={SCALE}x → {SLIDE_W*SCALE}×{SLIDE_H*SCALE}px)")

        paths = []
        for i, slide in enumerate(slides):
            out = SLIDES_DIR / f"slide_{i+1:02d}.png"
            slide.screenshot(path=str(out), type="png")
            paths.append(out)
            print(f"  Captured slide {i+1}/{len(slides)}")

        browser.close()
    return paths

def build_pptx(image_paths):
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W * 914400 // 96)   # 96 DPI → EMU
    prs.slide_height = Emu(SLIDE_H * 914400 // 96)
    blank_layout = prs.slide_layouts[6]  # blank

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path), Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )

    prs.save(str(OUTPUT_PATH))
    print(f"\nSaved → {OUTPUT_PATH}")

if __name__ == "__main__":
    images = capture_slides()
    build_pptx(images)
